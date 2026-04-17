from flask import Flask, request, jsonify, send_file, send_from_directory
from flask_cors import CORS
import os
import uuid
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

app = Flask(__name__, static_folder='../frontend', static_url_path='')
CORS(app)

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
OUTPUT_DIR = os.path.join(BASE_DIR, 'premium_ppts')
os.makedirs(OUTPUT_DIR, exist_ok=True)

CREATOR = "B. MARLIYA BANU"
APP_NAME = "PPT CREATOR"

# Professional color themes
COLOR_THEMES = {
    'corporate': {'primary': RGBColor(0, 51, 102), 'secondary': RGBColor(0, 112, 192), 'accent': RGBColor(255, 193, 7), 'name': 'Corporate Blue'},
    'modern': {'primary': RGBColor(75, 0, 130), 'secondary': RGBColor(138, 43, 226), 'accent': RGBColor(255, 140, 0), 'name': 'Modern Purple'},
    'minimal': {'primary': RGBColor(0, 80, 60), 'secondary': RGBColor(0, 120, 90), 'accent': RGBColor(255, 160, 0), 'name': 'Minimal Teal'},
    'bold': {'primary': RGBColor(180, 0, 0), 'secondary': RGBColor(220, 50, 50), 'accent': RGBColor(255, 215, 0), 'name': 'Bold Red'},
    'elegant': {'primary': RGBColor(100, 30, 80), 'secondary': RGBColor(156, 39, 176), 'accent': RGBColor(255, 193, 7), 'name': 'Elegant Purple'},
    'nature': {'primary': RGBColor(46, 125, 50), 'secondary': RGBColor(76, 175, 80), 'accent': RGBColor(255, 235, 59), 'name': 'Nature Green'},
    'ocean': {'primary': RGBColor(0, 150, 136), 'secondary': RGBColor(38, 198, 218), 'accent': RGBColor(255, 152, 0), 'name': 'Ocean Blue'}
}

def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs, title, creator_name, theme):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, RGBColor(255, 255, 255))
    
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.15))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = theme['primary']
    top_bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_box.text_frame.text = title
    title_box.text_frame.paragraphs[0].font.size = Pt(44)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.color.rgb = theme['primary']
    title_box.text_frame.paragraphs[0].font.name = 'Arial'
    title_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.6))
    sub_box.text_frame.text = "Professional Presentation"
    sub_box.text_frame.paragraphs[0].font.size = Pt(22)
    sub_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(80, 80, 80)
    sub_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    credit_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(0.5))
    credit_box.text_frame.text = f"Created by {creator_name}"
    credit_box.text_frame.paragraphs[0].font.size = Pt(16)
    credit_box.text_frame.paragraphs[0].font.bold = True
    credit_box.text_frame.paragraphs[0].font.color.rgb = theme['secondary']
    credit_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(0.4))
    date_box.text_frame.text = datetime.now().strftime('%B %d, %Y')
    date_box.text_frame.paragraphs[0].font.size = Pt(12)
    date_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(120, 120, 120)
    date_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_list, theme):
    """Standard content slide - fixed text display"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, RGBColor(250, 250, 250))
    
    # Title bar
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.9))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = theme['primary']
    title_bg.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.6))
    title_box.text_frame.text = title
    title_box.text_frame.paragraphs[0].font.size = Pt(28)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Content with proper spacing
    content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(8.8), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    # Calculate font size based on number of items
    num_items = len(content_list)
    if num_items <= 4:
        font_size = 24
    elif num_items <= 6:
        font_size = 20
    elif num_items <= 8:
        font_size = 18
    else:
        font_size = 16
    
    for i, line in enumerate(content_list):
        if line and str(line).strip():
            if i == 0:
                tf.text = f"▶ {str(line)}"
                tf.paragraphs[0].font.size = Pt(font_size)
                tf.paragraphs[0].font.color.rgb = RGBColor(40, 40, 40)
                tf.paragraphs[0].font.name = 'Calibri'
                tf.paragraphs[0].space_after = Pt(12)
            else:
                p = tf.add_paragraph()
                p.text = f"▶ {str(line)}"
                p.font.size = Pt(font_size)
                p.font.color.rgb = RGBColor(60, 60, 60)
                p.font.name = 'Calibri'
                p.space_after = Pt(10)

def add_roadmap_slide(prs, title, phases, theme):
    """Roadmap with connected round boxes - fixed text display"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, RGBColor(248, 248, 248))
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_box.text_frame.text = title
    title_box.text_frame.paragraphs[0].font.size = Pt(32)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.color.rgb = theme['primary']
    
    # Ensure at least 4 phases
    while len(phases) < 4:
        phases.append(f"Phase {len(phases)+1}")
    
    # Horizontal timeline line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.2), Inches(8.4), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = theme['primary']
    line.line.fill.background()
    
    positions = [1.2, 2.8, 4.4, 6.0, 7.6]
    
    for i, phase in enumerate(phases[:5]):
        x = positions[i]
        # Round circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x - 0.25), Inches(2.85), Inches(0.55), Inches(0.55))
        circle.fill.solid()
        circle.fill.fore_color.rgb = theme['accent']
        circle.line.fill.background()
        
        # Number inside circle
        num_box = slide.shapes.add_textbox(Inches(x - 0.18), Inches(2.93), Inches(0.45), Inches(0.4))
        num_box.text_frame.text = str(i + 1)
        num_box.text_frame.paragraphs[0].font.size = Pt(12)
        num_box.text_frame.paragraphs[0].font.bold = True
        num_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        num_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Phase box with proper text wrapping
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x - 0.8), Inches(1.5), Inches(1.6), Inches(1.0))
        box.fill.solid()
        box.fill.fore_color.rgb = theme['secondary']
        box.line.fill.background()
        box.text_frame.word_wrap = True
        
        # Handle text - split if too long
        phase_text = str(phase)
        if len(phase_text) > 15:
            # Try to split into two lines
            words = phase_text.split()
            if len(words) > 2:
                mid = len(words) // 2
                phase_text = words[0] + " " + words[1] + "\n" + " ".join(words[2:4])
            else:
                phase_text = phase_text[:12] + "\n" + phase_text[12:18] if len(phase_text) > 12 else phase_text
        
        box.text = phase_text
        box.text_frame.paragraphs[0].font.size = Pt(10)
        box.text_frame.paragraphs[0].font.bold = True
        box.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Arrow
        if i < len(phases[:5]) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 0.3), Inches(2.98), Inches(0.35), Inches(0.2))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = theme['primary']
            arrow.line.fill.background()

def add_mindmap_slide(prs, title, elements, theme):
    """Mind Map with proper center box and branches - fixed"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, RGBColor(250, 250, 250))
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_box.text_frame.text = title
    title_box.text_frame.paragraphs[0].font.size = Pt(32)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.color.rgb = theme['primary']
    
    # Center box - properly positioned
    center = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.2), Inches(2.5), Inches(3.2), Inches(1.0))
    center.fill.solid()
    center.fill.fore_color.rgb = theme['primary']
    center.line.color.rgb = theme['accent']
    center.line.width = Pt(2)
    center.text = "MAIN\nCONCEPT"
    center.text_frame.paragraphs[0].font.size = Pt(18)
    center.text_frame.paragraphs[0].font.bold = True
    center.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    center.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Ensure at least 4 elements
    while len(elements) < 4:
        elements.append(f"Point {len(elements)+1}")
    
    # Branch positions - adjusted for better visibility
    branches = [
        {"x": 0.6, "y": 1.5, "w": 1.8, "h": 0.8},
        {"x": 7.2, "y": 1.5, "w": 1.8, "h": 0.8},
        {"x": 0.6, "y": 4.2, "w": 1.8, "h": 0.8},
        {"x": 7.2, "y": 4.2, "w": 1.8, "h": 0.8}
    ]
    
    for i, (branch, elem) in enumerate(zip(branches, elements[:4])):
        # Connector line
        conn = slide.shapes.add_connector(1, Inches(4.8), Inches(3), Inches(branch['x'] + 0.9), Inches(branch['y'] + 0.4))
        conn.line.color.rgb = theme['secondary']
        conn.line.width = Pt(2)
        
        # Branch box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(branch['x']), Inches(branch['y']), Inches(branch['w']), Inches(branch['h']))
        box.fill.solid()
        box.fill.fore_color.rgb = theme['secondary']
        box.line.color.rgb = theme['accent']
        box.text_frame.word_wrap = True
        
        # Handle text - ensure it fits
        elem_text = str(elem)
        if len(elem_text) > 20:
            elem_text = elem_text[:18] + ".."
        box.text = elem_text
        box.text_frame.paragraphs[0].font.size = Pt(11)
        box.text_frame.paragraphs[0].font.bold = True
        box.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_features_slide(prs, title, features, theme):
    """Numbered features - fixed text display"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, RGBColor(248, 248, 248))
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_box.text_frame.text = title
    title_box.text_frame.paragraphs[0].font.size = Pt(32)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.color.rgb = theme['primary']
    
    # Ensure at least 4 features
    while len(features) < 4:
        features.append(f"Feature {len(features)+1}")
    
    colors = [theme['primary'], theme['secondary'], RGBColor(255, 140, 0), RGBColor(76, 175, 80), RGBColor(156, 39, 176)]
    
    for i, feature in enumerate(features[:6]):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.2 + i * 0.8), Inches(9), Inches(0.65))
        box.fill.solid()
        box.fill.fore_color.rgb = colors[i % len(colors)]
        box.line.fill.background()
        box.text_frame.word_wrap = True
        
        feature_text = str(feature)
        if len(feature_text) > 50:
            feature_text = feature_text[:47] + "..."
        
        box.text = f"{str(i+1).zfill(2)}  {feature_text}"
        box.text_frame.paragraphs[0].font.size = Pt(13)
        box.text_frame.paragraphs[0].font.bold = True
        box.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

def add_comparison_slide(prs, title, pros, cons, theme):
    """Pros vs Cons - fixed text display"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, RGBColor(250, 250, 250))
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_box.text_frame.text = title
    title_box.text_frame.paragraphs[0].font.size = Pt(32)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.color.rgb = theme['primary']
    
    # Ensure at least 3 pros and 3 cons
    while len(pros) < 3:
        pros.append(f"Benefit {len(pros)+1}")
    while len(cons) < 3:
        cons.append(f"Challenge {len(cons)+1}")
    
    # Pros Box
    pros_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.2), Inches(4.3), Inches(4.8))
    pros_box.fill.solid()
    pros_box.fill.fore_color.rgb = RGBColor(220, 255, 220)
    pros_box.line.color.rgb = RGBColor(0, 128, 0)
    pros_box.line.width = Pt(2)
    
    pros_title = pros_box.text_frame
    pros_title.text = "✓ PROS"
    pros_title.paragraphs[0].font.size = Pt(22)
    pros_title.paragraphs[0].font.bold = True
    pros_title.paragraphs[0].font.color.rgb = RGBColor(0, 128, 0)
    
    for item in pros[:6]:
        p = pros_title.add_paragraph()
        item_text = str(item)
        if len(item_text) > 35:
            item_text = item_text[:32] + "..."
        p.text = f"  • {item_text}"
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_after = Pt(8)
    
    # Cons Box
    cons_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.2), Inches(4.3), Inches(4.8))
    cons_box.fill.solid()
    cons_box.fill.fore_color.rgb = RGBColor(255, 220, 220)
    cons_box.line.color.rgb = RGBColor(255, 0, 0)
    cons_box.line.width = Pt(2)
    
    cons_title = cons_box.text_frame
    cons_title.text = "✗ CONS"
    cons_title.paragraphs[0].font.size = Pt(22)
    cons_title.paragraphs[0].font.bold = True
    cons_title.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)
    
    for item in cons[:6]:
        p = cons_title.add_paragraph()
        item_text = str(item)
        if len(item_text) > 35:
            item_text = item_text[:32] + "..."
        p.text = f"  • {item_text}"
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_after = Pt(8)

def add_conclusion_slide(prs, title, points, theme):
    """Conclusion Slide - fixed text display"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, theme['primary'])
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_box.text_frame.text = title
    title_box.text_frame.paragraphs[0].font.size = Pt(36)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    while len(points) < 4:
        points.append(f"Key Point {len(points)+1}")
    
    content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(8.8), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    font_size = 20 if len(points) <= 5 else 18
    
    for i, point in enumerate(points[:7]):
        point_text = str(point)
        if len(point_text) > 50:
            point_text = point_text[:47] + "..."
        
        if i == 0:
            tf.text = f"✓ {point_text}"
            tf.paragraphs[0].font.size = Pt(font_size)
            tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            tf.paragraphs[0].space_after = Pt(12)
        else:
            p = tf.add_paragraph()
            p.text = f"✓ {point_text}"
            p.font.size = Pt(font_size)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.space_after = Pt(10)

@app.route('/')
def index():
    return send_from_directory('../frontend', 'index.html')

@app.route('/<path:path>')
def serve_static(path):
    return send_from_directory('../frontend', path)

@app.route('/api/generate-ppt', methods=['POST'])
def generate_ppt():
    try:
        data = request.get_json()
        title = data.get('title', 'Presentation')
        creator_name = data.get('creator_name', CREATOR)
        style = data.get('style', 'corporate')
        slides_content = data.get('slides_content', [])
        
        theme = COLOR_THEMES.get(style, COLOR_THEMES['corporate'])
        prs = Presentation()
        
        # Title Slide
        add_title_slide(prs, title, creator_name, theme)
        
        # Template types
        template_types = ['content', 'roadmap', 'mindmap', 'features', 'comparison', 'conclusion']
        
        for i, slide in enumerate(slides_content[:14]):
            template_type = template_types[i % len(template_types)]
            slide_title = slide.get('title', f'Slide {i+1}')
            content = slide.get('content', [])
            
            # Clean content - remove empty strings
            content = [c for c in content if c and str(c).strip()]
            
            if template_type == 'content':
                add_content_slide(prs, slide_title, content, theme)
            elif template_type == 'roadmap':
                add_roadmap_slide(prs, slide_title, content, theme)
            elif template_type == 'mindmap':
                add_mindmap_slide(prs, slide_title, content, theme)
            elif template_type == 'features':
                add_features_slide(prs, slide_title, content, theme)
            elif template_type == 'comparison':
                mid = len(content) // 2 if len(content) > 0 else 3
                pros = content[:mid] if mid > 0 else ["Benefit 1", "Benefit 2", "Benefit 3"]
                cons = content[mid:] if mid < len(content) else ["Challenge 1", "Challenge 2", "Challenge 3"]
                add_comparison_slide(prs, slide_title, pros, cons, theme)
            elif template_type == 'conclusion':
                add_conclusion_slide(prs, slide_title, content, theme)
        
        filename = f"premium_ppt_{uuid.uuid4().hex[:6]}.pptx"
        filepath = os.path.join(OUTPUT_DIR, filename)
        prs.save(filepath)
        
        return jsonify({"download_url": f"/api/download/{filename}"})
    except Exception as e:
        print(f"Error: {e}")
        return jsonify({"error": str(e)}), 500

@app.route('/api/download/<filename>', methods=['GET'])
def download_ppt(filename):
    filepath = os.path.join(OUTPUT_DIR, filename)
    if os.path.exists(filepath):
        return send_file(filepath, as_attachment=True)
    return jsonify({"error": "File not found"}), 404

if __name__ == '__main__':
    print("\n" + "="*60)
    print(f"🐰 {APP_NAME} by {CREATOR}")
    print("="*60)
    print(f"✅ Server: http://127.0.0.1:5000")
    print("="*60)
    app.run(debug=True, port=5000)