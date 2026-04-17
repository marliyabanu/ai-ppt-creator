from flask import Flask, request, jsonify, send_file, send_from_directory
from flask_cors import CORS
import os
import uuid
from datetime import datetime
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

app = Flask(__name__, static_folder='../frontend', static_url_path='')
CORS(app)

OUTPUT_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'my_ppts')
os.makedirs(OUTPUT_DIR, exist_ok=True)

print(f"✅ Output folder: {OUTPUT_DIR}")

# Background colors
BACKGROUND_COLORS = {
    'white': RGBColor(255, 255, 255),
    'black': RGBColor(0, 0, 0),
    'light_gray': RGBColor(240, 240, 240),
    'dark_gray': RGBColor(50, 50, 50),
    'light_blue': RGBColor(230, 240, 255),
    'dark_blue': RGBColor(0, 51, 102),
    'light_green': RGBColor(230, 255, 230),
    'dark_green': RGBColor(0, 80, 60),
    'light_yellow': RGBColor(255, 255, 220),
    'light_purple': RGBColor(245, 230, 255),
    'dark_purple': RGBColor(75, 0, 130),
    'light_red': RGBColor(255, 230, 230),
    'dark_red': RGBColor(180, 0, 0)
}

# Font colors
FONT_COLORS = {
    'black': RGBColor(0, 0, 0),
    'white': RGBColor(255, 255, 255),
    'blue': RGBColor(0, 51, 102),
    'red': RGBColor(220, 20, 60),
    'green': RGBColor(0, 128, 0),
    'purple': RGBColor(128, 0, 128),
    'orange': RGBColor(255, 140, 0),
    'gray': RGBColor(80, 80, 80)
}

def get_font_size_by_content(content_list):
    """Calculate font size based on content length and text length"""
    # Count total characters in all content
    total_chars = sum(len(item) for item in content_list)
    num_items = len(content_list)
    
    # For long text, use smaller font
    if total_chars > 500:
        return 18
    elif total_chars > 400:
        return 20
    elif total_chars > 300:
        return 22
    elif total_chars > 200:
        return 24
    elif num_items <= 2:
        return 32
    elif num_items <= 3:
        return 28
    elif num_items <= 4:
        return 26
    elif num_items <= 5:
        return 24
    elif num_items <= 6:
        return 22
    else:
        return 20

def detect_content_type(title):
    """Detect content type from title"""
    t = title.lower()
    if 'advantage' in t or 'benefit' in t or 'pros' in t:
        return 'advantages'
    if 'disadvantage' in t or 'challenge' in t or 'cons' in t:
        return 'disadvantages'
    if 'feature' in t:
        return 'features'
    return 'standard'

def set_slide_background(slide, color):
    """Set slide background color"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

@app.route('/')
def index():
    return send_from_directory('../frontend', 'index.html')

@app.route('/<path:path>')
def serve_static(path):
    return send_from_directory('../frontend', path)

@app.route('/api/generate-ppt', methods=['POST'])
def generate_ppt():
    try:
        data = request.get_json()
        
        title = data.get('title', 'My Presentation')
        creator_name = data.get('creator_name', 'B. MARLIYA BANU')
        slides_content = data.get('slides_content', [])
        bg_color_name = data.get('background_color', 'white')
        font_color_name = data.get('font_color', 'black')
        font_style = data.get('font_style', 'Arial')
        
        bg_color = BACKGROUND_COLORS.get(bg_color_name, BACKGROUND_COLORS['white'])
        font_color = FONT_COLORS.get(font_color_name, FONT_COLORS['black'])
        
        print(f"Generating PPT: {title}")
        print(f"Slides: {len(slides_content)}")
        
        prs = Presentation()
        
        # ===== TITLE SLIDE =====
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        set_slide_background(slide, bg_color)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
        title_box.text_frame.text = title
        title_box.text_frame.paragraphs[0].font.size = Pt(44)
        title_box.text_frame.paragraphs[0].font.bold = True
        title_box.text_frame.paragraphs[0].font.color.rgb = font_color
        title_box.text_frame.paragraphs[0].font.name = font_style
        title_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Creator name
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
        sub_box.text_frame.text = f"Created by {creator_name}\n{datetime.now().strftime('%B %d, %Y')}"
        for para in sub_box.text_frame.paragraphs:
            para.font.size = Pt(18)
            para.font.color.rgb = font_color
            para.font.name = font_style
            para.alignment = PP_ALIGN.CENTER
        
        # ===== CONTENT SLIDES =====
        for idx, slide_data in enumerate(slides_content[:15]):
            slide_title = slide_data.get('title', f'Slide {idx+1}')
            content_raw = slide_data.get('content', [])
            
            # Clean content
            if isinstance(content_raw, str):
                content = [c.strip() for c in content_raw.split('\n') if c.strip()]
            else:
                content = [str(c).strip() for c in content_raw if c and str(c).strip()]
            
            if not content:
                content = ["No content provided"]
            
            content_type = detect_content_type(slide_title)
            font_size = get_font_size_by_content(content)
            
            print(f"Slide {idx+1}: '{slide_title}' - {len(content)} items - Font size: {font_size}pt")
            
            # Create slide
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
            set_slide_background(slide, bg_color)
            
            # Title
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
            title_box.text_frame.text = slide_title
            title_box.text_frame.paragraphs[0].font.size = Pt(32)
            title_box.text_frame.paragraphs[0].font.bold = True
            title_box.text_frame.paragraphs[0].font.color.rgb = font_color
            title_box.text_frame.paragraphs[0].font.name = font_style
            
            # Underline
            line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.0), Inches(2.5), Inches(0.05))
            line.fill.solid()
            line.fill.fore_color.rgb = font_color
            line.line.fill.background()
            
            # Content box
            content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(8.8), Inches(5.5))
            text_frame = content_box.text_frame
            text_frame.clear()
            text_frame.word_wrap = True
            
            if content_type == 'advantages':
                bullet_color = RGBColor(0, 128, 0)
                for item in content:
                    p = text_frame.add_paragraph()
                    p.text = f"✓  {item}"
                    p.font.size = Pt(font_size)
                    p.font.color.rgb = bullet_color
                    p.font.name = font_style
                    p.space_after = Pt(12)
                    p.line_spacing = 1.3
            elif content_type == 'disadvantages':
                bullet_color = RGBColor(220, 20, 60)
                for item in content:
                    p = text_frame.add_paragraph()
                    p.text = f"⚠  {item}"
                    p.font.size = Pt(font_size)
                    p.font.color.rgb = bullet_color
                    p.font.name = font_style
                    p.space_after = Pt(12)
                    p.line_spacing = 1.3
            elif content_type == 'features':
                for i, item in enumerate(content, 1):
                    p = text_frame.add_paragraph()
                    p.text = f"{i}.  {item}"
                    p.font.size = Pt(font_size)
                    p.font.color.rgb = font_color
                    p.font.name = font_style
                    p.space_after = Pt(12)
                    p.line_spacing = 1.3
            else:
                for item in content:
                    p = text_frame.add_paragraph()
                    p.text = f"•  {item}"
                    p.font.size = Pt(font_size)
                    p.font.color.rgb = font_color
                    p.font.name = font_style
                    p.space_after = Pt(12)
                    p.line_spacing = 1.3
        
        # Save file
        filename = f"presentation_{uuid.uuid4().hex[:6]}.pptx"
        filepath = os.path.join(OUTPUT_DIR, filename)
        prs.save(filepath)
        
        print(f"✅ PPT saved: {filepath}")
        
        return jsonify({"download_url": f"/api/download/{filename}"})
        
    except Exception as e:
        print(f"❌ Error: {str(e)}")
        import traceback
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500

@app.route('/api/download/<filename>', methods=['GET'])
def download_ppt(filename):
    filepath = os.path.join(OUTPUT_DIR, filename)
    if os.path.exists(filepath):
        return send_file(
            filepath,
            as_attachment=True,
            download_name=filename,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
    return jsonify({"error": "File not found"}), 404

if __name__ == '__main__':
    print("\n" + "="*60)
    print("🐰 PPT CREATOR by B. MARLIYA BANU")
    print("="*60)
    print("✅ Server: http://127.0.0.1:5000")
    print("✅ Auto font sizing - Long text = Smaller font")
    print("✅ Advantages: ✓ | Disadvantages: ⚠")
    print("="*60)
    app.run(debug=True, port=5000)