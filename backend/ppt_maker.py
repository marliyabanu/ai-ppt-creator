from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
import os
import uuid
from datetime import datetime

# Try to import pptx, if not available, create text file
try:
    from pptx import Presentation
    from pptx.util import Pt
    PPTX_AVAILABLE = True
    print("✅ python-pptx loaded successfully")
except ImportError:
    PPTX_AVAILABLE = False
    print("⚠️ python-pptx not installed. Will create text files instead.")
    print("📌 Install with: pip install python-pptx")

app = Flask(__name__)
CORS(app)

# Create output folder
OUTPUT_DIR = os.path.join(os.path.dirname(__file__), 'my_ppts')
os.makedirs(OUTPUT_DIR, exist_ok=True)
print(f"✅ Output folder: {OUTPUT_DIR}")

@app.route('/api/generate-outline', methods=['POST'])
def generate_outline():
    """Generate presentation outline"""
    try:
        data = request.get_json()
        prompt = data.get('prompt', 'Presentation')
        num_slides = int(data.get('num_slides', 20))
        
        outline = []
        
        # Title slide
        outline.append({
            "title": prompt,
            "content": [
                "Presented by AI PPT Creator",
                "Comprehensive Analysis",
                "Key Insights & Recommendations",
                "Strategic Overview"
            ]
        })
        
        # Complete presentation structure with Introduction, Features, Mind Map, Examples, Advantages, Disadvantages
        slides_content = [
            {
                "title": "Introduction",
                "content": [
                    f"Overview of {prompt}",
                    "Background and historical context",
                    "Why this topic is important today",
                    "Key objectives of this presentation",
                    "What you will learn"
                ]
            },
            {
                "title": "What is " + prompt.split()[0] + "?",
                "content": [
                    "Definition and core concepts",
                    "Key components and elements",
                    "How it works",
                    "Main principles and foundations",
                    "Terminology and key terms"
                ]
            },
            {
                "title": "Key Features & Characteristics",
                "content": [
                    "Primary features and capabilities",
                    "Unique selling points",
                    "Technical specifications",
                    "User benefits",
                    "Standout attributes"
                ]
            },
            {
                "title": "Mind Map & Conceptual Framework",
                "content": [
                    "Visual representation of core concepts",
                    "Main branches and sub-branches",
                    "Relationships and connections",
                    "Hierarchical structure",
                    "Key takeaways from mind map"
                ]
            },
            {
                "title": "Real-World Examples",
                "content": [
                    "Example 1: Industry application",
                    "Example 2: Business implementation",
                    "Example 3: Success story",
                    "Example 4: Innovative use case",
                    "Lessons learned from examples"
                ]
            },
            {
                "title": "Advantages & Benefits",
                "content": [
                    "Primary benefits and value proposition",
                    "Efficiency improvements",
                    "Cost savings and ROI",
                    "Competitive advantages",
                    "Long-term benefits"
                ]
            },
            {
                "title": "Disadvantages & Challenges",
                "content": [
                    "Limitations and constraints",
                    "Potential risks and issues",
                    "Implementation challenges",
                    "Common pitfalls to avoid",
                    "Mitigation strategies"
                ]
            },
            {
                "title": "Comparison: Pros vs Cons",
                "content": [
                    "Balanced analysis of strengths and weaknesses",
                    "When to use vs when to avoid",
                    "Risk vs reward assessment",
                    "Cost vs benefit analysis",
                    "Final verdict"
                ]
            },
            {
                "title": "Implementation Strategy",
                "content": [
                    "Step-by-step implementation plan",
                    "Resource requirements",
                    "Timeline and milestones",
                    "Team structure and roles",
                    "Success metrics and KPIs"
                ]
            },
            {
                "title": "Best Practices",
                "content": [
                    "Industry standards and guidelines",
                    "Proven strategies for success",
                    "Expert recommendations",
                    "Tips and tricks",
                    "Common mistakes to avoid"
                ]
            },
            {
                "title": "Case Study: Success Story",
                "content": [
                    "Company background",
                    "Challenge they faced",
                    "Solution implemented",
                    "Results achieved",
                    "Key takeaways"
                ]
            },
            {
                "title": "Future Trends & Opportunities",
                "content": [
                    "Emerging developments",
                    "Future predictions and forecasts",
                    "Growth opportunities",
                    "Innovation potential",
                    "What to watch for"
                ]
            },
            {
                "title": "Action Plan & Next Steps",
                "content": [
                    "Immediate actions (0-3 months)",
                    "Short-term goals (3-6 months)",
                    "Medium-term goals (6-12 months)",
                    "Long-term vision (1+ years)",
                    "Success measurement"
                ]
            },
            {
                "title": "Conclusion",
                "content": [
                    "Summary of key points",
                    "Main takeaways",
                    "Final recommendations",
                    "Call to action",
                    "Closing thoughts"
                ]
            },
            {
                "title": "Q&A Session",
                "content": [
                    "Discussion points",
                    "Feedback collection",
                    "Further clarification",
                    "Additional resources",
                    "Thank you"
                ]
            }
        ]
        
        # Add slides based on num_slides
        for i in range(min(num_slides - 1, len(slides_content))):
            outline.append(slides_content[i])
        
        # Add extra slides if needed
        while len(outline) < num_slides:
            outline.append({
                "title": f"Additional Insights - Part {len(outline)}",
                "content": [
                    "Deep dive into critical aspects",
                    "Supporting evidence and research",
                    "Practical applications",
                    "Future considerations",
                    "Key recommendations"
                ]
            })
        
        print(f"✅ Generated {len(outline)} slides for: {prompt}")
        return jsonify({"outline": outline})
        
    except Exception as e:
        print(f"❌ Error in generate_outline: {str(e)}")
        return jsonify({"error": str(e)}), 500

@app.route('/api/generate-ppt', methods=['POST'])
def generate_ppt():
    """Generate PowerPoint presentation"""
    try:
        data = request.get_json()
        outline = data.get('outline', [])
        style = data.get('style', 'professional')
        title = data.get('title', 'Presentation')
        
        # Generate unique filename
        filename = f"presentation_{uuid.uuid4().hex[:8]}.pptx"
        filepath = os.path.join(OUTPUT_DIR, filename)
        
        if PPTX_AVAILABLE:
            # Create real PowerPoint presentation
            prs = Presentation()
            
            # Style settings
            title_sizes = {'professional': 44, 'modern': 40, 'minimal': 48, 'bold': 52}
            text_sizes = {'professional': 24, 'modern': 22, 'minimal': 20, 'bold': 28}
            title_size = title_sizes.get(style, 44)
            text_size = text_sizes.get(style, 24)
            
            # Title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            slide.shapes.title.text = title
            slide.placeholders[1].text = f"Generated by AI PPT Creator\n{datetime.now().strftime('%B %d, %Y')}\nStyle: {style.upper()}"
            
            # Content slides
            for slide_data in outline:
                bullet_slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(bullet_slide_layout)
                
                # Set title
                slide.shapes.title.text = slide_data['title']
                for paragraph in slide.shapes.title.text_frame.paragraphs:
                    paragraph.font.size = Pt(title_size)
                    paragraph.font.bold = True
                
                # Set content
                content_shape = slide.placeholders[1]
                text_frame = content_shape.text_frame
                text_frame.clear()
                
                for point in slide_data['content']:
                    p = text_frame.add_paragraph()
                    p.text = point
                    p.font.size = Pt(text_size)
                    p.level = 0
                    p.space_after = Pt(12)
            
            # Save presentation
            prs.save(filepath)
            print(f"✅ Created PowerPoint: {filename}")
            
        else:
            # Create text file as fallback
            with open(filepath, 'w', encoding='utf-8') as f:
                f.write("="*70 + "\n")
                f.write(f"PRESENTATION: {title}\n")
                f.write(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
                f.write(f"Style: {style.upper()}\n")
                f.write(f"Total Slides: {len(outline)}\n")
                f.write("="*70 + "\n\n")
                
                for idx, slide in enumerate(outline, 1):
                    f.write(f"\n{'#'*70}\n")
                    f.write(f"SLIDE {idx}: {slide['title']}\n")
                    f.write(f"{'#'*70}\n")
                    for point in slide['content']:
                        f.write(f"  • {point}\n")
                    f.write("\n")
            print(f"✅ Created text file: {filename}")
        
        # Verify file exists
        if os.path.exists(filepath):
            file_size = os.path.getsize(filepath)
            print(f"📁 File size: {file_size} bytes")
            
            # Return download URL
            return jsonify({
                "success": True,
                "download_url": f"http://127.0.0.1:5000/api/download/{filename}",
                "filename": filename,
                "message": "Presentation generated successfully"
            })
        else:
            raise Exception("File was not created")
        
    except Exception as e:
        print(f"❌ Error in generate_ppt: {str(e)}")
        return jsonify({"error": str(e)}), 500

@app.route('/api/download/<filename>', methods=['GET'])
def download_ppt(filename):
    """Download the generated file"""
    try:
        # Security: prevent path traversal
        filename = os.path.basename(filename)
        filepath = os.path.join(OUTPUT_DIR, filename)
        
        print(f"📥 Download request: {filename}")
        print(f"📁 Looking for: {filepath}")
        
        if os.path.exists(filepath):
            file_size = os.path.getsize(filepath)
            print(f"✅ File found! Size: {file_size} bytes")
            
            # Determine mime type
            if filename.endswith('.pptx'):
                mimetype = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
            else:
                mimetype = 'application/octet-stream'
            
            return send_file(
                filepath,
                as_attachment=True,
                download_name=filename,
                mimetype=mimetype
            )
        else:
            print(f"❌ File not found: {filepath}")
            # List available files for debugging
            files = os.listdir(OUTPUT_DIR) if os.path.exists(OUTPUT_DIR) else []
            return jsonify({
                "error": f"File '{filename}' not found",
                "available_files": files,
                "output_directory": OUTPUT_DIR
            }), 404
            
    except Exception as e:
        print(f"❌ Download error: {str(e)}")
        return jsonify({"error": str(e)}), 500

@app.route('/api/list-files', methods=['GET'])
def list_files():
    """List all generated files"""
    try:
        files = os.listdir(OUTPUT_DIR) if os.path.exists(OUTPUT_DIR) else []
        return jsonify({
            "files": files,
            "count": len(files),
            "directory": OUTPUT_DIR
        })
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/api/health', methods=['GET'])
def health():
    """Health check endpoint"""
    files = os.listdir(OUTPUT_DIR) if os.path.exists(OUTPUT_DIR) else []
    return jsonify({
        "status": "healthy",
        "message": "Server is running!",
        "output_folder": OUTPUT_DIR,
        "files_count": len(files),
        "pptx_available": PPTX_AVAILABLE
    })

if __name__ == '__main__':
    print("\n" + "="*70)
    print("🎯 AI PPT CREATOR - BACKEND SERVER")
    print("="*70)
    print(f"✅ Server: http://127.0.0.1:5000")
    print(f"✅ Health check: http://127.0.0.1:5000/api/health")
    print(f"✅ Output folder: {OUTPUT_DIR}")
    print(f"✅ PPTX library: {'Available' if PPTX_AVAILABLE else 'Not available (using text fallback)'}")
    print("="*70)
    print("\n🚀 Ready to generate presentations!\n")
    app.run(debug=True, port=5000, host='127.0.0.1')