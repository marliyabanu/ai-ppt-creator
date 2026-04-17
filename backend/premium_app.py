from flask import Flask, request, jsonify, send_file, send_from_directory
from flask_cors import CORS
import os
import uuid
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

app = Flask(__name__, static_folder='../frontend', static_url_path='')
CORS(app)

OUTPUT_DIR = "premium_ppts"
os.makedirs(OUTPUT_DIR, exist_ok=True)

CREATOR_NAME = "B MARLIYA BANU"
APP_NAME = "PPT CREATOR"

def get_premium_theme(prompt):
    """Get premium theme based on prompt keywords"""
    prompt_lower = prompt.lower()
    
    themes = {
        'quantum': {
            'keywords': ['quantum', 'physics', 'atomic'],
            'primary': RGBColor(98, 0, 234), 'secondary': RGBColor(156, 39, 176),
            'accent': RGBColor(63, 81, 181), 'highlight': RGBColor(255, 193, 7),
            'title_font': 'Segoe UI', 'body_font': 'Calibri'
        },
        'ai': {
            'keywords': ['ai', 'artificial', 'intelligence', 'machine', 'neural', 'robot'],
            'primary': RGBColor(0, 150, 255), 'secondary': RGBColor(0, 200, 255),
            'accent': RGBColor(0, 100, 200), 'highlight': RGBColor(255, 100, 0),
            'title_font': 'Segoe UI', 'body_font': 'Calibri'
        },
        'business': {
            'keywords': ['business', 'finance', 'corporate', 'strategy', 'management', 'marketing'],
            'primary': RGBColor(0, 100, 80), 'secondary': RGBColor(0, 150, 120),
            'accent': RGBColor(255, 140, 0), 'highlight': RGBColor(255, 215, 0),
            'title_font': 'Arial', 'body_font': 'Calibri'
        },
        'health': {
            'keywords': ['health', 'medical', 'wellness', 'fitness', 'care'],
            'primary': RGBColor(220, 20, 60), 'secondary': RGBColor(255, 100, 150),
            'accent': RGBColor(128, 0, 0), 'highlight': RGBColor(255, 200, 200),
            'title_font': 'Georgia', 'body_font': 'Calibri'
        },
        'environment': {
            'keywords': ['environment', 'climate', 'nature', 'green', 'sustainable', 'earth'],
            'primary': RGBColor(0, 120, 0), 'secondary': RGBColor(100, 200, 100),
            'accent': RGBColor(60, 80, 20), 'highlight': RGBColor(255, 200, 100),
            'title_font': 'Century Gothic', 'body_font': 'Calibri'
        },
        'technology': {
            'keywords': ['tech', 'digital', 'software', 'coding', 'blockchain', 'cyber', 'data'],
            'primary': RGBColor(0, 112, 192), 'secondary': RGBColor(0, 176, 240),
            'accent': RGBColor(0, 32, 96), 'highlight': RGBColor(255, 140, 0),
            'title_font': 'Consolas', 'body_font': 'Segoe UI'
        }
    }
    
    default = {
        'primary': RGBColor(0, 0, 0), 'secondary': RGBColor(80, 80, 80),
        'accent': RGBColor(50, 50, 50), 'highlight': RGBColor(200, 200, 200),
        'title_font': 'Arial', 'body_font': 'Arial'
    }
    
    for theme in themes.values():
        if any(kw in prompt_lower for kw in theme['keywords']):
            return theme
    return default

def add_title_slide(prs, title, theme):
    """Add branded title slide"""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(48)
    slide.shapes.title.text_frame.paragraphs[0].font.bold = True
    slide.shapes.title.text_frame.paragraphs[0].font.name = theme['title_font']
    
    # Subtitle with branding
    subtitle = slide.placeholders[1]
    subtitle.text = f"Presented by {APP_NAME}\nCreated by {CREATOR_NAME}\n{datetime.now().strftime('%B %d, %Y')}"
    for para in subtitle.text_frame.paragraphs:
        para.font.size = Pt(18)
        para.font.name = theme['body_font']

def add_mind_map_slide(prs, title, theme):
    """Add premium mind map with boxes and arrows"""
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    title_box.text_frame.text = title
    title_box.text_frame.paragraphs[0].font.size = Pt(36)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.name = theme['title_font']
    
    # Center box
    center = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(2.5), Inches(2.5), Inches(1))
    center.fill.solid()
    center.fill.fore_color.rgb = theme['primary']
    center.line.color.rgb = theme['accent']
    center.text = "MAIN CONCEPT"
    center.text_frame.paragraphs[0].font.size = Pt(18)
    center.text_frame.paragraphs[0].font.bold = True
    center.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Branch boxes
    branches = [
        {"name": "INTRODUCTION", "x": 0.5, "y": 1.5},
        {"name": "FEATURES", "x": 7, "y": 1.5},
        {"name": "BENEFITS", "x": 0.5, "y": 4.5},
        {"name": "APPLICATIONS", "x": 7, "y": 4.5}
    ]
    
    for b in branches:
        # Connector line
        conn = slide.shapes.add_connector(1, Inches(4.75), Inches(3), Inches(b['x'] + 0.8), Inches(b['y'] + 0.3))
        conn.line.color.rgb = theme['secondary']
        conn.line.width = Pt(2)
        
        # Branch box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(b['x']), Inches(b['y']), Inches(1.6), Inches(0.7))
        box.fill.solid()
        box.fill.fore_color.rgb = theme['secondary']
        box.text = b['name']
        box.text_frame.paragraphs[0].font.size = Pt(12)
        box.text_frame.paragraphs[0].font.bold = True

def add_comparison_slide(prs, title, theme):
    """Add advantages vs disadvantages slide"""
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    title_box.text_frame.text = title
    title_box.text_frame.paragraphs[0].font.size = Pt(36)
    title_box.text_frame.paragraphs[0].font.bold = True
    
    # Advantages box (Left)
    adv_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.2), Inches(4), Inches(4.5))
    adv_box.fill.solid()
    adv_box.fill.fore_color.rgb = RGBColor(220, 255, 220)
    adv_box.line.color.rgb = RGBColor(0, 128, 0)
    adv_box.line.width = Pt(2)
    
    adv_title = adv_box.text_frame
    adv_title.text = "✓ ADVANTAGES"
    adv_title.paragraphs[0].font.size = Pt(24)
    adv_title.paragraphs[0].font.bold = True
    adv_title.paragraphs[0].font.color.rgb = RGBColor(0, 128, 0)
    
    # Disadvantages box (Right)
    dis_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.5), Inches(1.2), Inches(4), Inches(4.5))
    dis_box.fill.solid()
    dis_box.fill.fore_color.rgb = RGBColor(255, 220, 220)
    dis_box.line.color.rgb = RGBColor(255, 0, 0)
    dis_box.line.width = Pt(2)
    
    dis_title = dis_box.text_frame
    dis_title.text = "✗ DISADVANTAGES"
    dis_title.paragraphs[0].font.size = Pt(24)
    dis_title.paragraphs[0].font.bold = True
    dis_title.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)

def add_feature_slide(prs, title, theme):
    """Add premium feature slide with icons"""
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    title_box.text_frame.text = title
    title_box.text_frame.paragraphs[0].font.size = Pt(36)
    title_box.text_frame.paragraphs[0].font.bold = True
    
    # Feature boxes
    features = [
        {"name": "FEATURE 1", "desc": "Innovative solution", "x": 0.5, "y": 1.5},
        {"name": "FEATURE 2", "desc": "User friendly", "x": 3.5, "y": 1.5},
        {"name": "FEATURE 3", "desc": "Cost effective", "x": 6.5, "y": 1.5},
        {"name": "FEATURE 4", "desc": "Scalable", "x": 0.5, "y": 3.5},
        {"name": "FEATURE 5", "desc": "Secure", "x": 3.5, "y": 3.5},
        {"name": "FEATURE 6", "desc": "Reliable", "x": 6.5, "y": 3.5}
    ]
    
    for f in features:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(f['x']), Inches(f['y']), Inches(2.5), Inches(1.2))
        box.fill.solid()
        box.fill.fore_color.rgb = theme['primary']
        box.text = f"{f['name']}\n{f['desc']}"
        box.text_frame.paragraphs[0].font.size = Pt(14)
        box.text_frame.paragraphs[0].font.bold = True
        box.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

@app.route('/')
def index():
    return send_from_directory('../frontend', 'index.html')

@app.route('/<path:path>')
def serve_static(path):
    return send_from_directory('../frontend', path)

@app.route('/api/generate-outline', methods=['POST'])
def generate_outline():
    try:
        data = request.get_json()
        prompt = data.get('prompt', 'Presentation')
        num_slides = int(data.get('num_slides', 15))
        
        outline = []
        
        # Premium slide structure
        slides = [
            {"title": "Introduction", "type": "normal", "content": [f"Overview of {prompt}", "Why this matters", "Key objectives", "What you'll learn"]},
            {"title": "What is " + prompt.split()[0] + "?", "type": "normal", "content": ["Definition and core concepts", "Historical background", "Key components", "How it works"]},
            {"title": "Key Features", "type": "features", "content": []},
            {"title": "Mind Map", "type": "mindmap", "content": []},
            {"title": "Real-World Examples", "type": "normal", "content": ["Example 1: Industry application", "Example 2: Business success", "Example 3: Innovative use case"]},
            {"title": "Advantages vs Disadvantages", "type": "comparison", "content": []},
            {"title": "Benefits & Opportunities", "type": "normal", "content": ["Benefit 1: Efficiency improvement", "Benefit 2: Cost savings", "Benefit 3: Competitive advantage", "Benefit 4: Growth potential"]},
            {"title": "Challenges & Solutions", "type": "normal", "content": ["Challenge 1: Implementation", "Solution: Strategic planning", "Challenge 2: Cost", "Solution: Phased approach"]},
            {"title": "Future Trends", "type": "normal", "content": ["Trend 1: Emerging technologies", "Trend 2: Market growth", "Trend 3: Innovation opportunities"]},
            {"title": "Conclusion", "type": "normal", "content": ["Summary of key points", "Main takeaways", "Call to action", "Next steps"]}
        ]
        
        for i in range(min(num_slides, len(slides))):
            outline.append(slides[i])
        
        while len(outline) < num_slides:
            outline.append({"title": f"Additional Insights {len(outline)}", "type": "normal", "content": ["Key point 1", "Key point 2", "Key point 3"]})
        
        return jsonify({"outline": outline})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/api/generate-ppt', methods=['POST'])
def generate_ppt():
    try:
        data = request.get_json()
        outline = data.get('outline', [])
        style = data.get('style', 'professional')
        title = data.get('title', 'Presentation')
        
        theme = get_premium_theme(title)
        prs = Presentation()
        
        # Add title slide with branding
        add_title_slide(prs, title, theme)
        
        # Add content slides
        for slide_data in outline:
            slide_type = slide_data.get('type', 'normal')
            
            if slide_type == 'mindmap':
                add_mind_map_slide(prs, slide_data['title'], theme)
            elif slide_type == 'comparison':
                add_comparison_slide(prs, slide_data['title'], theme)
            elif slide_type == 'features':
                add_feature_slide(prs, slide_data['title'], theme)
            else:
                # Normal bullet slide
                bullet_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(bullet_layout)
                slide.shapes.title.text = slide_data['title']
                
                content = slide.placeholders[1]
                text_frame = content.text_frame
                text_frame.clear()
                
                for point in slide_data.get('content', []):
                    p = text_frame.add_paragraph()
                    p.text = point
                    p.font.size = Pt(20)
                    p.font.name = theme['body_font']
                    p.space_after = Pt(10)
        
        # Save file
        filename = f"premium_ppt_{uuid.uuid4().hex[:6]}.pptx"
        filepath = os.path.join(OUTPUT_DIR, filename)
        prs.save(filepath)
        
        return jsonify({"download_url": f"/api/download/{filename}"})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/api/download/<filename>', methods=['GET'])
def download_ppt(filename):
    filepath = os.path.join(OUTPUT_DIR, filename)
    if os.path.exists(filepath):
        return send_file(filepath, as_attachment=True)
    return jsonify({"error": "File not found"}), 404

if __name__ == '__main__':
    print("\n" + "="*60)
    print(f"🌟 {APP_NAME} by {CREATOR_NAME}")
    print("="*60)
    print(f"✅ Server: http://127.0.0.1:5000")
    print("="*60)
    app.run(debug=True, port=5000)